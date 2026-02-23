#!/usr/bin/env python3
"""
Extract text content from PowerPoint (.pptx) files.

Usage:
    python3 extract_pptx.py <path_to_pptx_file>
"""

import sys
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    """Extract all text from a PowerPoint presentation."""
    try:
        prs = Presentation(pptx_path)

        print(f"Total slides: {len(prs.slides)}\n")
        print("=" * 80)

        for slide_num, slide in enumerate(prs.slides, start=1):
            print(f"\n## SLIDE {slide_num}")
            print("-" * 80)

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"\n{shape.text}")

                # Handle tables
                if shape.has_table:
                    print("\n[TABLE]")
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        print(row_text)
                    print("[/TABLE]\n")

            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    print(f"\n**Notes:** {notes_text}")

            print("\n" + "=" * 80)

        return True

    except Exception as e:
        print(f"Error reading PowerPoint file: {e}", file=sys.stderr)
        return False

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python3 extract_pptx.py <path_to_pptx_file>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    success = extract_text_from_pptx(pptx_path)
    sys.exit(0 if success else 1)
