#!/usr/bin/env python3
"""
Extract reusable templates from PowerPoint presentations.

Analyses a .pptx file and extracts its design system (slide masters, layouts,
colour themes, font themes, placeholder configurations) into a reusable
template catalogue stored as JSON.

Usage:
    python3 pptx_extract_template.py <path_to_pptx> [--name "Template Name"]
    python3 pptx_extract_template.py <path_to_pptx> --dry-run
    python3 pptx_extract_template.py --list
"""

import argparse
import json
import logging
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.util import Emu, Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from lxml import etree
except ImportError:
    print("Error: lxml is required. Install with: pip3 install lxml", file=sys.stderr)
    sys.exit(1)

logging.basicConfig(level=logging.WARNING, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

# Resolve paths relative to script location
SCRIPT_DIR = Path(__file__).resolve().parent
VAULT_ROOT = SCRIPT_DIR.parent.parent
TEMPLATES_DIR = VAULT_ROOT / ".claude" / "templates" / "pptx"
CATALOGUE_PATH = TEMPLATES_DIR / "_catalogue.json"

DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": DRAWINGML_NS}


def slugify(name: str) -> str:
    """Convert a name to a kebab-case slug."""
    slug = name.lower().strip()
    slug = re.sub(r"[^a-z0-9\s-]", "", slug)
    slug = re.sub(r"[\s_]+", "-", slug)
    slug = re.sub(r"-+", "-", slug)
    return slug.strip("-")


def emu_to_inches(emu: int) -> float:
    """Convert EMU to inches for human-readable display."""
    return round(emu / 914400, 2) if emu else 0.0


def extract_colour(element) -> str | None:
    """Extract hex colour from a DrawingML colour element."""
    if element is None:
        return None

    # srgbClr — explicit RGB
    srgb = element.find(f"{{{DRAWINGML_NS}}}srgbClr")
    if srgb is not None:
        return f"#{srgb.get('val', '000000').upper()}"

    # sysClr — system colour with lastClr fallback
    sys_clr = element.find(f"{{{DRAWINGML_NS}}}sysClr")
    if sys_clr is not None:
        last = sys_clr.get("lastClr")
        if last:
            return f"#{last.upper()}"
        return f"sys:{sys_clr.get('val', 'unknown')}"

    return None


def get_theme_element(master):
    """Get the theme XML element from a slide master's related parts."""
    try:
        theme_part = master.part.part_related_by(RT.THEME)
        return etree.fromstring(theme_part.blob)
    except Exception:
        logger.debug("Could not access theme part via relationship, falling back to element search")
        return master.element


def extract_colour_scheme(master) -> dict:
    """Extract the colour scheme from a slide master's theme."""
    colours = {}
    theme_elem = get_theme_element(master)

    clr_scheme = theme_elem.find(f".//{{{DRAWINGML_NS}}}clrScheme")
    if clr_scheme is None:
        logger.debug("No colour scheme found in theme")
        return colours

    colour_names = [
        "dk1", "dk2", "lt1", "lt2",
        "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
        "hlink", "folHlink",
    ]
    for name in colour_names:
        elem = clr_scheme.find(f"{{{DRAWINGML_NS}}}{name}")
        colour = extract_colour(elem)
        if colour:
            colours[name] = colour

    return colours


def extract_font_scheme(master) -> dict:
    """Extract the font scheme from a slide master's theme."""
    fonts = {"major": None, "minor": None}
    theme_elem = get_theme_element(master)

    font_scheme = theme_elem.find(f".//{{{DRAWINGML_NS}}}fontScheme")
    if font_scheme is None:
        logger.debug("No font scheme found in theme")
        return fonts

    major = font_scheme.find(f".//{{{DRAWINGML_NS}}}majorFont/{{{DRAWINGML_NS}}}latin")
    minor = font_scheme.find(f".//{{{DRAWINGML_NS}}}minorFont/{{{DRAWINGML_NS}}}latin")

    if major is not None:
        fonts["major"] = major.get("typeface")
    if minor is not None:
        fonts["minor"] = minor.get("typeface")

    return fonts


def extract_background_info(background) -> dict | None:
    """Extract background fill information from a slide/master background."""
    if background is None:
        return None

    bg_info = {}
    fill = background.fill

    try:
        fill_type = fill.type
        if fill_type is not None:
            bg_info["fillType"] = str(fill_type)
        else:
            bg_info["fillType"] = "inherited"
    except Exception:
        bg_info["fillType"] = "unknown"

    return bg_info if bg_info else None


def extract_placeholder_info(placeholder) -> dict:
    """Extract detailed information about a placeholder."""
    ph_info = {
        "idx": placeholder.placeholder_format.idx,
        "type": str(placeholder.placeholder_format.type),
        "left": placeholder.left,
        "top": placeholder.top,
        "width": placeholder.width,
        "height": placeholder.height,
        "leftInches": emu_to_inches(placeholder.left),
        "topInches": emu_to_inches(placeholder.top),
        "widthInches": emu_to_inches(placeholder.width),
        "heightInches": emu_to_inches(placeholder.height),
    }

    # Extract text frame defaults if present
    if placeholder.has_text_frame:
        tf = placeholder.text_frame
        text_info = {}

        if tf.paragraphs:
            first_para = tf.paragraphs[0]
            if first_para.alignment is not None:
                text_info["alignment"] = str(first_para.alignment)

            font = first_para.font
            if font.name:
                text_info["fontName"] = font.name
            if font.size:
                text_info["fontSize"] = font.size.pt
            if font.bold is not None:
                text_info["bold"] = font.bold
            if font.italic is not None:
                text_info["italic"] = font.italic

        if text_info:
            ph_info["textDefaults"] = text_info

    return ph_info


def extract_layout_info(layout, layout_index: int, master_index: int) -> dict:
    """Extract detailed information about a slide layout."""
    layout_info = {
        "name": layout.name or f"Layout {layout_index}",
        "slug": slugify(layout.name or f"layout-{layout_index}"),
        "index": layout_index,
        "masterIndex": master_index,
        "placeholders": [],
    }

    # Extract placeholder details
    for ph in layout.placeholders:
        ph_info = extract_placeholder_info(ph)
        layout_info["placeholders"].append(ph_info)

    # Sort placeholders by idx for consistency
    layout_info["placeholders"].sort(key=lambda p: p["idx"])

    return layout_info


def extract_master_info(master, master_index: int) -> dict:
    """Extract information about a slide master."""
    master_info = {
        "index": master_index,
        "layoutCount": len(master.slide_layouts),
        "placeholderCount": len(master.placeholders),
    }

    bg = extract_background_info(master.background)
    if bg:
        master_info["background"] = bg

    return master_info


def deduplicate_layouts(layouts: list[dict]) -> list[dict]:
    """Remove duplicate and empty layouts, keeping the first occurrence of each unique layout.

    A layout is considered a duplicate if it has the same name AND the same
    placeholder signature (set of placeholder types and indices).
    Empty 'Default Slide' layouts with no placeholders are filtered out.
    """
    seen = set()
    unique = []

    for layout in layouts:
        name = layout["name"]
        placeholders = layout.get("placeholders", [])

        # Skip empty default slides
        if name == "Default Slide" and not placeholders:
            continue

        # Build a signature: name + sorted placeholder idx:type pairs
        ph_sig = tuple(
            sorted((p["idx"], p["type"]) for p in placeholders)
        )
        key = (name, ph_sig)

        if key not in seen:
            seen.add(key)
            unique.append(layout)
        else:
            logger.debug(f"Skipping duplicate layout: {name}")

    return unique


def extract_template(pptx_path: Path) -> dict:
    """Extract the full template metadata from a PPTX file."""
    logger.info(f"Opening presentation: {pptx_path}")
    prs = Presentation(str(pptx_path))

    template = {
        "slideWidth": prs.slide_width,
        "slideHeight": prs.slide_height,
        "slideWidthInches": emu_to_inches(prs.slide_width),
        "slideHeightInches": emu_to_inches(prs.slide_height),
        "slideCount": len(prs.slides),
        "slideMasters": [],
        "layouts": [],
        "theme": {
            "colours": {},
            "fonts": {"major": None, "minor": None},
        },
    }

    # Extract slide masters
    all_layouts = []
    for m_idx, master in enumerate(prs.slide_masters):
        master_info = extract_master_info(master, m_idx)
        template["slideMasters"].append(master_info)

        # Extract theme from first master (primary theme)
        if m_idx == 0:
            template["theme"]["colours"] = extract_colour_scheme(master)
            template["theme"]["fonts"] = extract_font_scheme(master)

        # Extract layouts under this master
        for l_idx, layout in enumerate(master.slide_layouts):
            layout_info = extract_layout_info(layout, l_idx, m_idx)
            all_layouts.append(layout_info)

    # Deduplicate layouts
    raw_count = len(all_layouts)
    template["layouts"] = deduplicate_layouts(all_layouts)
    dedup_count = len(template["layouts"])
    if raw_count != dedup_count:
        logger.info(f"Deduplicated layouts: {raw_count} -> {dedup_count}")

    return template


def generate_previews(
    pptx_path: Path, template_dir: Path, layouts: list[dict]
) -> list[str]:
    """Generate PNG preview thumbnails for each layout using LibreOffice.

    Creates a temporary PPTX with one blank slide per layout, converts to
    PDF via LibreOffice, then splits into per-layout PNGs.
    """
    previews_dir = template_dir / "previews"
    previews_dir.mkdir(exist_ok=True)

    soffice = shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")

    if not soffice or not pdftoppm:
        logger.warning("LibreOffice or poppler not found — skipping preview generation")
        return []

    generated = []

    try:
        # Create a temp PPTX with one slide per layout
        prs = Presentation(str(pptx_path))

        # Remove existing slides
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            r_id = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if r_id:
                prs.part.drop_rel(r_id)
            sld_id_lst.remove(sld_id)

        # Add one slide per layout
        for layout_data in layouts:
            m_idx = layout_data.get("masterIndex", 0)
            l_idx = layout_data.get("index", 0)
            try:
                layout = prs.slide_masters[m_idx].slide_layouts[l_idx]
                prs.slides.add_slide(layout)
            except (IndexError, KeyError):
                logger.debug(f"Could not add slide for layout {layout_data['name']}")

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_pptx = Path(tmpdir) / "layouts.pptx"
            prs.save(str(tmp_pptx))

            # Convert to PDF
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, str(tmp_pptx)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                return []

            tmp_pdf = Path(tmpdir) / "layouts.pdf"
            if not tmp_pdf.exists():
                logger.warning("PDF not created by LibreOffice")
                return []

            # Convert PDF pages to PNGs
            result = subprocess.run(
                [pdftoppm, "-png", "-r", "150", str(tmp_pdf), str(Path(tmpdir) / "slide")],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                logger.warning(f"pdftoppm conversion failed: {result.stderr}")
                return []

            # Move generated PNGs to previews directory
            png_files = sorted(Path(tmpdir).glob("slide-*.png"))
            for i, png_file in enumerate(png_files):
                if i < len(layouts):
                    slug = layouts[i]["slug"]
                    dest = previews_dir / f"{slug}.png"
                    shutil.copy2(png_file, dest)
                    generated.append(f"{slug}.png")

    except Exception as e:
        logger.warning(f"Preview generation failed: {e}")

    return generated


def save_template(
    pptx_path: Path,
    template_data: dict,
    name: str,
    slug: str,
    tags: list[str] | None = None,
    generate_preview: bool = False,
) -> Path:
    """Save extracted template to the catalogue directory."""
    template_dir = TEMPLATES_DIR / slug
    template_dir.mkdir(parents=True, exist_ok=True)

    # Copy source PPTX
    source_dest = template_dir / "source.pptx"
    shutil.copy2(pptx_path, source_dest)
    logger.info(f"Copied source PPTX to {source_dest}")

    # Build full template metadata
    full_template = {
        "name": name,
        "slug": slug,
        "source": pptx_path.name,
        "extracted": date.today().isoformat(),
        **template_data,
    }

    # Save template.json
    template_json_path = template_dir / "template.json"
    with open(template_json_path, "w", encoding="utf-8") as f:
        json.dump(full_template, f, indent=2, ensure_ascii=False)
    logger.info(f"Saved template metadata to {template_json_path}")

    # Save individual layout files
    layouts_dir = template_dir / "layouts"
    layouts_dir.mkdir(exist_ok=True)
    for layout in template_data.get("layouts", []):
        layout_path = layouts_dir / f"{layout['slug']}.json"
        with open(layout_path, "w", encoding="utf-8") as f:
            json.dump(layout, f, indent=2, ensure_ascii=False)

    # Save theme files
    theme_dir = template_dir / "theme"
    theme_dir.mkdir(exist_ok=True)
    theme = template_data.get("theme", {})
    with open(theme_dir / "colours.json", "w", encoding="utf-8") as f:
        json.dump(theme.get("colours", {}), f, indent=2, ensure_ascii=False)
    with open(theme_dir / "fonts.json", "w", encoding="utf-8") as f:
        json.dump(theme.get("fonts", {}), f, indent=2, ensure_ascii=False)

    # Generate previews if requested
    previews = []
    if generate_preview:
        previews = generate_previews(pptx_path, template_dir, template_data.get("layouts", []))

    # Update catalogue
    update_catalogue(name, slug, pptx_path.name, template_data, tags or [], previews)

    return template_dir


def update_catalogue(
    name: str,
    slug: str,
    source_filename: str,
    template_data: dict,
    tags: list[str],
    previews: list[str],
) -> None:
    """Update the master catalogue file."""
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    catalogue = {"version": 1, "templates": []}
    if CATALOGUE_PATH.exists():
        with open(CATALOGUE_PATH, "r", encoding="utf-8") as f:
            catalogue = json.load(f)

    # Remove existing entry with same slug (idempotent)
    catalogue["templates"] = [t for t in catalogue["templates"] if t.get("slug") != slug]

    layout_names = [l["name"] for l in template_data.get("layouts", [])]
    entry = {
        "name": name,
        "slug": slug,
        "source": source_filename,
        "extracted": date.today().isoformat(),
        "layoutCount": len(template_data.get("layouts", [])),
        "layoutNames": layout_names,
        "tags": tags,
        "hasPreviews": len(previews) > 0,
        "slideCount": template_data.get("slideCount", 0),
        "theme": {
            "colours": list(template_data.get("theme", {}).get("colours", {}).values())[:6],
            "majorFont": template_data.get("theme", {}).get("fonts", {}).get("major"),
            "minorFont": template_data.get("theme", {}).get("fonts", {}).get("minor"),
        },
    }

    catalogue["templates"].append(entry)

    with open(CATALOGUE_PATH, "w", encoding="utf-8") as f:
        json.dump(catalogue, f, indent=2, ensure_ascii=False)

    logger.info(f"Updated catalogue at {CATALOGUE_PATH}")


def print_summary(name: str, slug: str, template_data: dict, template_dir: Path | None = None) -> None:
    """Print a human-readable extraction summary."""
    layouts = template_data.get("layouts", [])
    theme = template_data.get("theme", {})
    colours = theme.get("colours", {})
    fonts = theme.get("fonts", {})

    print(f"\n{'=' * 60}")
    print(f"  Template Extraction Summary")
    print(f"{'=' * 60}")
    print(f"  Name:         {name}")
    print(f"  Slug:         {slug}")
    print(f"  Slide size:   {template_data.get('slideWidthInches', '?')}\" x {template_data.get('slideHeightInches', '?')}\"")
    print(f"  Source slides: {template_data.get('slideCount', '?')}")
    print(f"  Masters:      {len(template_data.get('slideMasters', []))}")
    print(f"  Layouts:      {len(layouts)}")
    print()

    # Layouts table
    print(f"  {'Layout Name':<30} {'Placeholders':>13}")
    print(f"  {'-' * 30} {'-' * 13}")
    for layout in layouts:
        ph_count = len(layout.get("placeholders", []))
        ph_types = ", ".join(
            p.get("type", "?").replace("PP_PLACEHOLDER_TYPE.", "").replace("PP_PLACEHOLDER.", "")
            for p in layout.get("placeholders", [])
        )
        print(f"  {layout['name']:<30} {ph_count:>5}   ({ph_types})")

    print()

    # Theme
    if colours:
        print(f"  Theme Colours:")
        for k, v in list(colours.items())[:8]:
            print(f"    {k:<12} {v}")
    if fonts.get("major") or fonts.get("minor"):
        print(f"\n  Theme Fonts:")
        print(f"    Major (headings): {fonts.get('major', 'not set')}")
        print(f"    Minor (body):     {fonts.get('minor', 'not set')}")

    if template_dir:
        print(f"\n  Saved to: {template_dir}")

    print(f"{'=' * 60}\n")


def list_templates() -> None:
    """List all templates in the catalogue."""
    if not CATALOGUE_PATH.exists():
        print("No template catalogue found. Extract a template first.")
        return

    with open(CATALOGUE_PATH, "r", encoding="utf-8") as f:
        catalogue = json.load(f)

    templates = catalogue.get("templates", [])
    if not templates:
        print("No templates in catalogue.")
        return

    print(f"\n{'=' * 60}")
    print(f"  PPTX Template Catalogue ({len(templates)} templates)")
    print(f"{'=' * 60}")

    for t in templates:
        print(f"\n  {t['name']}")
        print(f"    Slug:     {t['slug']}")
        print(f"    Source:   {t['source']}")
        print(f"    Layouts:  {t.get('layoutCount', '?')}")
        print(f"    Extracted: {t.get('extracted', '?')}")
        if t.get("tags"):
            print(f"    Tags:     {', '.join(t['tags'])}")

    print(f"\n{'=' * 60}\n")


def main():
    parser = argparse.ArgumentParser(
        description="Extract reusable templates from PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pptx
  %(prog)s presentation.pptx --name "BA Community Day"
  %(prog)s presentation.pptx --dry-run
  %(prog)s presentation.pptx --previews
  %(prog)s --list
        """,
    )
    parser.add_argument("pptx_path", nargs="?", help="Path to the PowerPoint file")
    parser.add_argument("--name", help="Template name (defaults to filename)")
    parser.add_argument("--tags", nargs="*", default=[], help="Tags for the template")
    parser.add_argument("--dry-run", action="store_true", help="Analyse only, do not save")
    parser.add_argument("--previews", action="store_true", help="Generate layout preview images")
    parser.add_argument("--list", action="store_true", help="List all extracted templates")
    parser.add_argument("--verbose", "-v", action="store_true", help="Enable debug logging")

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    if args.list:
        list_templates()
        return

    if not args.pptx_path:
        parser.error("pptx_path is required unless --list is used")

    pptx_path = Path(args.pptx_path).expanduser().resolve()
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if not pptx_path.suffix.lower() in (".pptx", ".ppt"):
        print(f"Warning: File does not have a .pptx extension: {pptx_path}", file=sys.stderr)

    # Determine template name and slug
    name = args.name or pptx_path.stem.replace("_", " ").replace("-", " ")
    slug = slugify(name)

    print(f"Extracting template from: {pptx_path.name}")
    print(f"Template name: {name} (slug: {slug})")
    print()

    # Extract
    template_data = extract_template(pptx_path)

    # Print summary
    if args.dry_run:
        print("[DRY RUN] Analysis complete — no files saved.\n")
        print_summary(name, slug, template_data)
    else:
        template_dir = save_template(
            pptx_path,
            template_data,
            name,
            slug,
            tags=args.tags,
            generate_preview=args.previews,
        )
        print_summary(name, slug, template_data, template_dir)


if __name__ == "__main__":
    main()
